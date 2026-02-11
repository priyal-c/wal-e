"""
PPTX deck reporter - Generates WAL_Assessment_Presentation.pptx using python-pptx.

Produces a 16-slide executive readout matching the original WAL assessment deck exactly:
  1.  Title
  2.  Assessment Objective
  3.  Overall Maturity Radar (big score + maturity legend table)
  4.  Pillar Summary Scores (score bars ■□ + maturity per pillar)
  5.  Workspace at a Glance (inventory + key config, dual-column)
  6.  Top 10 Critical Findings (with Quick Fix? column)
  7-13. Pillar Deep Dives (Strengths / Critical Gaps / Key Recommendations numbered)
  14. Remediation Roadmap (2x2 grid)
  15. Recommended Next Steps (immediate + follow-up engagements table)
  16. Assessment Summary Statistics (BP breakdown table)
  17. Thank You & Contact (deliverables + references)
"""

from pathlib import Path
from typing import Any, Dict, List, Optional, Union

from .base import (
    PILLAR_DISPLAY_NAMES,
    PILLAR_ORDER,
    AuditEntry,
    BaseReporter,
    BestPracticeScore,
    ScoredAssessment,
)

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    Presentation = None  # type: ignore

# ---------------------------------------------------------------------------
# Databricks-inspired color palette
# ---------------------------------------------------------------------------
if PPTX_AVAILABLE:
    DB_RED = RGBColor(0xFF, 0x38, 0x21)
    DB_DARK = RGBColor(0x1B, 0x1F, 0x23)
    DB_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    DB_LIGHT = RGBColor(0xE8, 0xEA, 0xED)
    DB_GRAY = RGBColor(0x9A, 0xA0, 0xA6)
    DB_GREEN = RGBColor(0x00, 0xA9, 0x72)
    DB_ORANGE = RGBColor(0xF5, 0x9E, 0x0B)
    DB_CRIT = RGBColor(0xEF, 0x44, 0x44)
    DB_BLUE = RGBColor(0x13, 0x6C, 0xB9)
    TBL_HDR = RGBColor(0x2D, 0x33, 0x3A)
    TBL_ODD = RGBColor(0x22, 0x27, 0x2E)
    TBL_EVEN = RGBColor(0x1B, 0x1F, 0x23)
else:
    DB_RED = DB_DARK = DB_WHITE = DB_LIGHT = DB_GRAY = None
    DB_GREEN = DB_ORANGE = DB_CRIT = DB_BLUE = None
    TBL_HDR = TBL_ODD = TBL_EVEN = None


class PPTXDeckReporter(BaseReporter):
    """Generates a polished PPTX executive readout presentation."""

    def __init__(self):
        super().__init__("WAL_Assessment_Presentation.pptx")

    def generate(self, scored_assessment: ScoredAssessment, collected_data: Dict[str, Any],
                 audit_entries: List[AuditEntry], output_dir: Union[str, Path]) -> Path:
        if not PPTX_AVAILABLE:
            raise ImportError("python-pptx is not installed. pip install 'python-pptx>=1.0'")

        out = self._ensure_output_dir(output_dir) / self.output_filename
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        ps = self._get_pillar_scores(scored_assessment)
        bps = self._get_best_practice_scores(scored_assessment)
        ov = self._get_overall_score(scored_assessment)
        mat = self._get_maturity_level(scored_assessment)
        host = self._get_workspace_host(scored_assessment)
        date = self._get_assessment_date(scored_assessment)

        gov = collected_data.get("GovernanceCollector", {})
        comp = collected_data.get("ComputeCollector", {})
        ops = collected_data.get("OperationsCollector", {})
        sec = collected_data.get("SecurityCollector", {})
        ws = collected_data.get("WorkspaceCollector", {})

        self._s1_title(prs, host, date, gov)
        self._s2_objective(prs)
        self._s3_maturity_radar(prs, ps, ov, mat)
        self._s4_pillar_scores(prs, ps)
        self._s5_workspace(prs, gov, comp, ops, sec)
        self._s6_findings(prs, bps)
        for pillar in PILLAR_ORDER:
            self._s_pillar(prs, pillar, ps, bps, collected_data)
        self._s_roadmap(prs, bps, sec, comp)
        self._s_next_steps(prs, sec, comp, gov)
        self._s_stats(prs, bps, ps)
        self._s_thankyou(prs, host)

        prs.save(str(out))
        return out

    # ====================== HELPERS ======================

    def _bg(self, slide):
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = DB_DARK

    def _title(self, slide, text, top=0.4, sz=32):
        tx = slide.shapes.add_textbox(Inches(0.8), Inches(top), Inches(11.5), Inches(0.8))
        p = tx.text_frame.paragraphs[0]
        p.text = text; p.font.size = Pt(sz); p.font.bold = True; p.font.color.rgb = DB_RED

    def _txt(self, slide, text, left, top, width=11.5, sz=12, color=None, bold=False):
        tx = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(0.4))
        tx.text_frame.word_wrap = True
        p = tx.text_frame.paragraphs[0]
        p.text = text; p.font.size = Pt(sz); p.font.color.rgb = color or DB_LIGHT; p.font.bold = bold

    def _bullets(self, slide, items, left, top, width=5.5, sz=11, color=None):
        y = top
        for item in items:
            tx = slide.shapes.add_textbox(Inches(left), Inches(y), Inches(width), Inches(0.35))
            tx.text_frame.word_wrap = True
            p = tx.text_frame.paragraphs[0]
            p.text = f"\u2022  {item}"; p.font.size = Pt(sz); p.font.color.rgb = color or DB_LIGHT
            y += 0.35
        return y

    def _numbered(self, slide, items, left, top, width=5.5, sz=11, color=None):
        y = top
        for i, item in enumerate(items, 1):
            tx = slide.shapes.add_textbox(Inches(left), Inches(y), Inches(width), Inches(0.35))
            tx.text_frame.word_wrap = True
            p = tx.text_frame.paragraphs[0]
            p.text = f"{i}. {item}"; p.font.size = Pt(sz); p.font.color.rgb = color or DB_LIGHT
            y += 0.35
        return y

    def _cell(self, cell, sz, fg, bg, bold=False):
        cell.fill.solid(); cell.fill.fore_color.rgb = bg
        for p in cell.text_frame.paragraphs:
            p.font.size = sz; p.font.color.rgb = fg; p.font.bold = bold
        cell.margin_left = Inches(0.08); cell.margin_right = Inches(0.08)
        cell.margin_top = Inches(0.04); cell.margin_bottom = Inches(0.04)

    def _health_color(self, health: str):
        h = health.upper()
        if "CRITICAL" in h: return DB_CRIT
        if "WARNING" in h: return DB_ORANGE
        return DB_GREEN

    def _risk_color(self, risk: str):
        r = risk.upper()
        if "CRITICAL" in r: return DB_CRIT
        if "HIGH" in r: return DB_ORANGE
        return DB_GREEN

    def _score_color(self, s):
        if s >= 1.5: return DB_GREEN
        if s >= 0.5: return DB_ORANGE
        return DB_CRIT

    def _mat(self, s):
        if s >= 1.75: return "Optimized"
        if s >= 1.25: return "Established"
        if s >= 0.5: return "Developing"
        return "Beginning"

    def _bar(self, score, width=10):
        pct = min(1.0, max(0, score / 2.0))
        filled = int(width * pct)
        return "\u25A0" * filled + "\u25A1" * (width - filled)

    def _tbl(self, slide, rows, hdrs, left, top, width, cw=None, color_cols=None):
        """Add styled table. color_cols: dict mapping col_idx -> 'health'|'risk'|'score' for conditional coloring."""
        nr = len(rows) + 1; nc = len(hdrs)
        sh = slide.shapes.add_table(nr, nc, Inches(left), Inches(top), Inches(width), Inches(0.3 * nr))
        tbl = sh.table
        if cw:
            for i, w in enumerate(cw):
                tbl.columns[i].width = Inches(w)
        for j, h in enumerate(hdrs):
            c = tbl.cell(0, j); c.text = h
            self._cell(c, Pt(10), DB_WHITE, TBL_HDR, bold=True)
        for i, row in enumerate(rows):
            bg = TBL_ODD if i % 2 == 0 else TBL_EVEN
            for j, val in enumerate(row):
                c = tbl.cell(i + 1, j); c.text = str(val)
                fg = DB_LIGHT
                if color_cols and j in color_cols:
                    kind = color_cols[j]
                    if kind == "health": fg = self._health_color(val)
                    elif kind == "risk": fg = self._risk_color(val)
                    elif kind == "score":
                        try: fg = self._score_color(float(val))
                        except ValueError: pass
                    elif kind == "maturity":
                        v = val.lower()
                        if "optimized" in v or "established" in v: fg = DB_GREEN
                        elif "developing" in v: fg = DB_ORANGE
                        elif "beginning" in v: fg = DB_CRIT
                    elif kind == "quickfix":
                        fg = DB_GREEN if val.lower() == "yes" else DB_ORANGE
                self._cell(c, Pt(9), fg, bg)

    # ====================== SLIDE 1: TITLE ======================

    def _s1_title(self, prs, host, date, gov):
        s = prs.slides.add_slide(prs.slide_layouts[6]); self._bg(s)
        tx = s.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.5), Inches(1.5))
        p = tx.text_frame.paragraphs[0]
        p.text = "Well-Architected Lakehouse Assessment"; p.font.size = Pt(44); p.font.bold = True; p.font.color.rgb = DB_WHITE
        p2 = tx.text_frame.add_paragraph()
        p2.text = "Executive Readout"; p2.font.size = Pt(44); p2.font.bold = True; p2.font.color.rgb = DB_RED

        meta = [
            ("Assessment Type", "Well-Architected Lakehouse Assessment"),
            ("Customer Workspace", host),
            ("Metastore", gov.get("metastore_name", "N/A")),
            ("Assessment Date", date),
            ("Lead Assessor", "WAL-E Agent (Automated)"),
            ("Framework", "Databricks Well-Architected Framework"),
        ]
        self._tbl(s, [[l, v] for l, v in meta], ["Field", "Details"], 0.8, 3.8, 10, cw=[3.0, 7.0])

    # ====================== SLIDE 2: OBJECTIVE ======================

    def _s2_objective(self, prs):
        s = prs.slides.add_slide(prs.slide_layouts[6]); self._bg(s)
        self._title(s, "Assessment Objective")
        self._txt(s, "Objective: Evaluate the existing Databricks workspace architecture against industry "
                  "standard best practices and provide high-level recommendations based on the Well-Architected "
                  "Lakehouse framework.", 0.8, 1.4, sz=14, bold=True)
        self._txt(s, "Scope: All 7 pillars of the WAL framework assessed:", 0.8, 2.2, sz=14, bold=True)
        pillars = ["Data & AI Governance", "Interoperability & Usability", "Operational Excellence",
                   "Security, Compliance & Privacy", "Reliability", "Performance Efficiency", "Cost Optimization"]
        self._numbered(s, pillars, 1.0, 2.7, sz=12)
        self._txt(s, "Methodology: Automated workspace inspection via Databricks CLI, REST APIs, Unity Catalog APIs, "
                  "and workspace configuration analysis.", 0.8, 5.4, sz=14, bold=True)

    # ====================== SLIDE 3: OVERALL MATURITY ======================

    def _s3_maturity_radar(self, prs, ps, ov, mat):
        s = prs.slides.add_slide(prs.slide_layouts[6]); self._bg(s)
        self._title(s, "Overall Maturity Summary")

        # Big score
        tx = s.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(4), Inches(1.2))
        p = tx.text_frame.paragraphs[0]
        pct = self._score_to_pct(ov)
        p.text = f"{pct:.0f}%"; p.font.size = Pt(64); p.font.bold = True; p.font.color.rgb = self._score_color(ov)
        p2 = tx.text_frame.add_paragraph()
        p2.text = f"Overall Score: {ov:.2f} / 2.0 ({mat} Maturity)"
        p2.font.size = Pt(14); p2.font.color.rgb = DB_GRAY

        # Maturity legend table (right side)
        legend = [
            ["Beginning", "0.0 - 0.5", "Minimal best practices adopted; significant gaps"],
            ["Developing", "0.5 - 1.25", "Some best practices in place; inconsistent adoption"],
            ["Established", "1.25 - 1.75", "Most best practices adopted; room for optimization"],
            ["Optimized", "1.75 - 2.0", "Full best practice adoption; continuous improvement"],
        ]
        self._tbl(s, legend, ["Maturity Level", "Score Range", "Description"],
                  5.5, 1.3, 7.3, cw=[2.0, 1.5, 3.8], color_cols={0: "maturity"})

        # Pillar mini-scores below
        y = 3.8
        for pillar in PILLAR_ORDER:
            sc = ps.get(pillar, 0)
            pct_p = self._score_to_pct(sc)
            disp = self._pillar_display_name(pillar)
            bar = self._bar(sc)
            self._txt(s, f"{disp}", 0.8, y, width=5.0, sz=12, color=DB_LIGHT)
            self._txt(s, f"{bar}  {sc:.1f}  ({pct_p:.0f}%)", 6.5, y, width=6.0, sz=12, color=self._score_color(sc))
            y += 0.42

    # ====================== SLIDE 4: PILLAR SUMMARY SCORES ======================

    def _s4_pillar_scores(self, prs, ps):
        s = prs.slides.add_slide(prs.slide_layouts[6]); self._bg(s)
        self._title(s, "Pillar Summary Scores")

        rows = []
        for pillar in PILLAR_ORDER:
            sc = ps.get(pillar, 0)
            disp = self._pillar_display_name(pillar)
            bar = self._bar(sc)
            m = self._mat(sc)
            rows.append([disp, f"{sc:.1f}", bar, m])

        self._tbl(s, rows, ["Pillar", "Score", "Visual", "Maturity"],
                  0.5, 1.3, 12.3, cw=[4.5, 1.2, 3.5, 3.1],
                  color_cols={1: "score", 3: "maturity"})

    # ====================== SLIDE 5: WORKSPACE AT A GLANCE ======================

    def _s5_workspace(self, prs, gov, comp, ops, sec):
        s = prs.slides.add_slide(prs.slide_layouts[6]); self._bg(s)
        self._title(s, "Workspace at a Glance")
        self._txt(s, "Resource Inventory", 0.5, 1.2, sz=14, bold=True)

        cc = gov.get("catalog_count", 0)
        el = gov.get("external_location_count", 0)
        sc_ = gov.get("storage_credential_count", 0)
        pc = comp.get("policy_count", 0)
        rc = comp.get("running_clusters", comp.get("cluster_count", 0))
        wc = comp.get("warehouse_count", 0)
        pp = comp.get("pool_count", 0)
        pl = ops.get("pipeline_count", 0)
        ep = ops.get("endpoint_count", 0)
        rp = ops.get("repo_count", 0)
        sc_cnt = ops.get("scope_count", 0)

        inv = [
            ["Unity Catalog Catalogs", str(cc), "CRITICAL - Extreme sprawl" if cc > 100 else ("WARNING - Review" if cc > 20 else "OK")],
            ["External Locations", str(el), "WARNING - Significant sprawl" if el > 50 else "OK"],
            ["Storage Credentials", str(sc_), "WARNING - Significant sprawl" if sc_ > 50 else "OK"],
            ["Cluster Policies", str(pc), "CRITICAL - No standardization" if pc == 0 else ("WARNING - Too many" if pc > 50 else "OK")],
            ["Running Clusters", str(rc), "CRITICAL - High cost exposure" if rc > 10 else "OK"],
            ["SQL Warehouses", str(wc), "WARNING - Excessive count" if wc > 20 else "OK"],
            ["Instance Pools", str(pp), "WARNING - Review utilization" if pp > 30 else "OK"],
            ["DLT Pipelines", str(pl), "OK - Active" if pl > 0 else "WARNING - None"],
            ["Serving Endpoints", str(ep), "OK - Active and serving" if ep > 0 else "OK"],
            ["Secret Scopes", str(sc_cnt), "OK" if sc_cnt > 0 else "WARNING - None"],
            ["Git Repos", str(rp), "OK - Active integration" if rp > 0 else "WARNING - None"],
        ]
        self._tbl(s, inv, ["Resource", "Count", "Health"], 0.3, 1.6, 6.2, cw=[2.6, 1.0, 2.6], color_cols={2: "health"})

        # Key Configuration (right column)
        ss = sec.get("security_settings", {}) or {}
        self._txt(s, "Key Configuration", 7.0, 1.2, sz=14, bold=True)
        conf = []
        dbfs = ss.get("enableDbfsFileBrowser", "")
        conf.append(["DBFS File Browser", str(dbfs).upper() if dbfs else "N/A", "DISABLE" if str(dbfs).lower() == "true" else "MAINTAIN"])
        dl = ss.get("enableResultsDownloading", "")
        conf.append(["Results Download", str(dl).upper() if dl else "N/A", "EVALUATE"])
        exp = ss.get("enableExportNotebook", "")
        conf.append(["Notebook Export", str(exp).upper() if exp else "N/A", "EVALUATE"])
        ipl = ss.get("enableIpAccessLists", "")
        conf.append(["IP Access Lists", str(ipl).upper() if ipl else "N/A", "MAINTAIN" if str(ipl).lower() == "true" else "ENABLE"])
        mtl = ss.get("maxTokenLifetimeDays", "")
        conf.append(["Max Token Lifetime", f"{mtl} days" if mtl else "N/A", "30 days" if mtl and int(mtl) > 30 else "OK"])
        ms = gov.get("metastore_summary", {})
        if isinstance(ms, dict):
            owner = ms.get("owner", "N/A")
            conf.append(["Metastore Owner", str(owner), "Admin group" if "account" in str(owner).lower() else "OK"])
        iso = gov.get("isolation_modes", [])
        if isinstance(iso, list) and "OPEN" in iso:
            conf.append(["Catalog Isolation", "Mostly OPEN", "ISOLATED"])
        self._tbl(s, conf, ["Setting", "Current Value", "Recommended"], 6.8, 1.6, 6.2, cw=[2.4, 1.8, 2.0])

    # ====================== SLIDE 6: TOP 10 FINDINGS ======================

    def _s6_findings(self, prs, bps):
        s = prs.slides.add_slide(prs.slide_layouts[6]); self._bg(s)
        self._title(s, "Top 10 Critical Findings")

        top = sorted([b for b in bps if b.get("score") is not None],
                     key=lambda x: (float(x.get("score", 2)), x.get("pillar", "")))[:10]

        rows = []
        for i, bp in enumerate(top, 1):
            sc = float(bp.get("score", 0))
            name = bp.get("name", "")
            notes = bp.get("finding_notes", "")
            finding = f"{name} - {notes}"[:75]
            disp = self._pillar_display_name(bp.get("pillar", ""))
            risk = "CRITICAL" if sc == 0 else ("HIGH" if sc <= 1 else "MEDIUM")
            # Quick fix: score-0 items with simple config changes are quick fixes
            quick = "Yes" if sc == 0 and any(kw in notes.lower() for kw in ["disable", "owner", "configure", "enable", "reduce", "define"]) else "No"
            rows.append([str(i), finding, disp, risk, quick])

        self._tbl(s, rows, ["#", "Finding", "Pillar", "Risk", "Quick Fix?"],
                  0.3, 1.3, 12.7, cw=[0.4, 5.8, 3.0, 1.5, 2.0],
                  color_cols={3: "risk", 4: "quickfix"})

    # ====================== SLIDES 7-13: PILLAR DEEP DIVES ======================

    def _s_pillar(self, prs, pillar, ps, bps, cd):
        s = prs.slides.add_slide(prs.slide_layouts[6]); self._bg(s)
        disp = self._pillar_display_name(pillar)
        sc = ps.get(pillar, 0)
        pct = self._score_to_pct(sc)
        m = self._mat(sc)
        self._title(s, f"{disp} Deep Dive")
        self._txt(s, f"Score: {sc:.1f}/2.0 - {m}", 0.8, 1.1, sz=16, bold=True, color=self._score_color(sc))

        pbps = self._get_bps_for_pillar(bps, pillar)
        strengths = [b for b in pbps if float(b.get("score", 0)) == 2]
        gaps = [b for b in pbps if float(b.get("score", 0)) == 0]
        partial = [b for b in pbps if float(b.get("score", 0)) == 1]

        # LEFT COLUMN: Strengths
        y = 1.7
        self._txt(s, "Strengths:", 0.8, y, width=5.5, sz=13, bold=True, color=DB_GREEN)
        y += 0.35
        if strengths:
            items = [f"{b.get('finding_notes', b.get('name', ''))}"[:65] for b in strengths[:5]]
        else:
            items = ["No fully implemented best practices in this pillar"]
        y = self._bullets(s, items, 0.8, y, width=5.5, sz=10, color=DB_LIGHT)
        y += 0.15

        # Critical Gaps
        self._txt(s, "Critical Gaps:", 0.8, y, width=5.5, sz=13, bold=True, color=DB_CRIT)
        y += 0.35
        if gaps:
            items = [f"{b.get('finding_notes', b.get('name', ''))}"[:65] for b in gaps[:6]]
        else:
            items = ["No critical gaps (score 0) in this pillar"]
        y = self._bullets(s, items, 0.8, y, width=5.5, sz=10, color=DB_LIGHT)

        # RIGHT COLUMN: Key Recommendations (numbered)
        self._txt(s, "Key Recommendations:", 7.0, 1.7, width=5.5, sz=13, bold=True, color=DB_WHITE)
        recs = []
        for b in gaps:
            n = b.get("finding_notes", "")
            if n: recs.append(n[:65])
        for b in partial[:5]:
            n = b.get("finding_notes", "")
            if n and len(recs) < 6: recs.append(n[:65])
        if not recs:
            recs = ["Maintain current best practices"]
        self._numbered(s, recs[:6], 7.0, 2.1, width=5.5, sz=10, color=DB_LIGHT)

        # Bottom: compact BP score table (if fits)
        if len(pbps) <= 15:
            tbl_top = max(y + 0.3, 5.0)
            row_h = 0.25
            if tbl_top + row_h * (len(pbps) + 1) <= 7.3:
                rows = []
                for b in pbps:
                    sv = float(b.get("score", 0))
                    rows.append([b.get("name", ""), f"{sv:.0f}", b.get("finding_notes", "")[:50]])
                self._tbl(s, rows, ["Best Practice", "Score", "Notes"],
                         0.3, tbl_top, 12.7, cw=[3.2, 0.8, 8.7], color_cols={1: "score"})

    # ====================== SLIDE 14: ROADMAP (2x2 GRID) ======================

    def _s_roadmap(self, prs, bps, sec, comp):
        s = prs.slides.add_slide(prs.slide_layouts[6]); self._bg(s)
        self._title(s, "Remediation Roadmap Summary")

        ss = sec.get("security_settings", {}) or {}
        zero = [b for b in bps if float(b.get("score", 2)) == 0]
        partial = [b for b in bps if float(b.get("score", 2)) == 1]

        # Phase 1 - top left
        self._txt(s, "PHASE 1 (Week 1-2)", 0.5, 1.3, width=5.5, sz=14, bold=True, color=DB_RED)
        self._txt(s, "Quick Wins", 0.5, 1.7, width=5.5, sz=12, bold=True, color=DB_GRAY)
        p1 = []
        if str(ss.get("enableDbfsFileBrowser", "")).lower() == "true":
            p1.append("Disable DBFS browser")
        if zero:
            for b in zero[:3]:
                p1.append(b.get("name", "")[:40])
        mtl = ss.get("maxTokenLifetimeDays", "")
        if mtl and int(mtl) > 30:
            p1.append("Reduce token lifetime")
        if comp.get("running_clusters", 0) > 10:
            p1.append("Terminate idle clusters")
        self._bullets(s, p1[:6], 0.5, 2.1, width=5.5, sz=10)

        # Phase 2 - top right
        self._txt(s, "PHASE 2 (Week 3-6)", 7.0, 1.3, width=5.5, sz=14, bold=True, color=DB_ORANGE)
        self._txt(s, "Foundation", 7.0, 1.7, width=5.5, sz=12, bold=True, color=DB_GRAY)
        p2 = ["Catalog strategy & naming convention", "Standardize cluster policies", "Cost tagging strategy",
              "Catalog isolation (OPEN -> ISOLATED)", "Restrict IAM profiles", "Audit logging via system tables"]
        self._bullets(s, p2[:6], 7.0, 2.1, width=5.5, sz=10)

        # Phase 3 - bottom left
        self._txt(s, "PHASE 3 (Week 7-12)", 0.5, 4.3, width=5.5, sz=14, bold=True, color=DB_BLUE)
        self._txt(s, "Operational Maturity", 0.5, 4.7, width=5.5, sz=12, bold=True, color=DB_GRAY)
        p3 = ["Dev/Test/Prod separation", "Standardize CI/CD (DAB/Terraform)", "Monitoring dashboards",
              "Fix failing DLT pipelines", "DQ expectations framework", "Consolidate SQL Warehouses"]
        self._bullets(s, p3[:6], 0.5, 5.1, width=5.5, sz=10)

        # Phase 4 - bottom right
        self._txt(s, "PHASE 4 (Week 13-20)", 7.0, 4.3, width=5.5, sz=14, bold=True, color=DB_GREEN)
        self._txt(s, "Optimization", 7.0, 4.7, width=5.5, sz=12, bold=True, color=DB_GRAY)
        p4 = ["Disaster recovery plan", "Private Link / VPC injection", "Liquid clustering / Z-ORDER",
              "Chargeback / showback system", "Performance testing in CI/CD", "Deploy Security Analysis Tool (SAT)"]
        self._bullets(s, p4[:6], 7.0, 5.1, width=5.5, sz=10)

    # ====================== SLIDE 15: NEXT STEPS ======================

    def _s_next_steps(self, prs, sec, comp, gov):
        s = prs.slides.add_slide(prs.slide_layouts[6]); self._bg(s)
        self._title(s, "Recommended Next Steps")

        ss = sec.get("security_settings", {}) or {}
        self._txt(s, "Immediate Actions (This Week)", 0.8, 1.2, sz=14, bold=True)

        imm = []
        if str(ss.get("enableDbfsFileBrowser", "")).lower() == "true":
            imm.append("Disable DBFS file browser - Workspace admin setting change (5 minutes)")
        ms = gov.get("metastore_summary", {})
        if isinstance(ms, dict) and "account" in str(ms.get("owner", "")).lower():
            imm.append("Change metastore owner from 'account users' to dedicated admin group (5 minutes)")
        mtl = ss.get("maxTokenLifetimeDays", "")
        if mtl and int(mtl) > 30:
            imm.append(f"Reduce token lifetime from {mtl} to 30 days (5 minutes)")
        if comp.get("running_clusters", 0) > 10:
            imm.append("Terminate idle running clusters - Review and stop non-essential clusters")
        if comp.get("warehouse_count", 0) > 20:
            imm.append(f"Stop unused SQL Warehouses - Audit {comp.get('warehouse_count', 0)} warehouses for actual usage")
        if not imm:
            imm = ["Prioritize findings based on business impact", "Assign owners for Phase 1-4 initiatives",
                   "Schedule follow-up assessment in 90 days"]
        self._numbered(s, imm[:5], 0.8, 1.7, width=11.5, sz=12)

        # Follow-up engagements
        y = 1.7 + 0.35 * min(5, len(imm)) + 0.4
        self._txt(s, "Follow-Up Engagements", 0.8, y, sz=14, bold=True)
        eng = [
            ["Catalog Cleanup & Governance", "Define catalog strategy, naming conventions, lifecycle policies", "2-3 weeks"],
            ["Security Hardening", "IAM profile review, Private Link, SAT deployment, compliance framework", "3-4 weeks"],
            ["Platform Standardization", "Cluster policy templates, CI/CD standardization, monitoring setup", "4-6 weeks"],
            ["Cost Optimization Workshop", "Tagging strategy, chargeback implementation, resource consolidation", "2-3 weeks"],
            ["Disaster Recovery Planning", "DR design, cross-region replication, backup strategy", "3-4 weeks"],
        ]
        self._tbl(s, eng, ["Engagement", "Description", "Duration"], 0.3, y + 0.4, 12.7, cw=[3.5, 6.7, 2.5])

    # ====================== SLIDE 16: STATS ======================

    def _s_stats(self, prs, bps, ps):
        s = prs.slides.add_slide(prs.slide_layouts[6]); self._bg(s)
        self._title(s, "Assessment Tool Summary")
        self._txt(s, "Scored Best Practices by Pillar", 0.8, 1.2, sz=14, bold=True)

        rows = []; t_bp = t_0 = t_1 = t_2 = 0
        for pillar in PILLAR_ORDER:
            pbps = self._get_bps_for_pillar(bps, pillar)
            n = len(pbps)
            s0 = sum(1 for b in pbps if float(b.get("score", 0)) == 0)
            s1 = sum(1 for b in pbps if float(b.get("score", 0)) == 1)
            s2 = sum(1 for b in pbps if float(b.get("score", 0)) == 2)
            avg = ps.get(pillar, 0)
            d = self._pillar_display_name(pillar)
            rows.append([d, str(n), str(s0), str(s1), str(s2), f"{avg:.2f}"])
            t_bp += n; t_0 += s0; t_1 += s1; t_2 += s2

        t_avg = sum(ps.values()) / len(ps) if ps else 0
        rows.append(["TOTAL", str(t_bp), str(t_0), str(t_1), str(t_2), f"{t_avg:.2f}"])

        self._tbl(s, rows, ["Pillar", "Total BPs Assessed", "Score 0", "Score 1", "Score 2", "Avg Score"],
                  0.3, 1.6, 12.7, cw=[4.0, 2.0, 1.3, 1.3, 1.3, 2.8], color_cols={5: "score"})

        pct0 = (t_0 / t_bp * 100) if t_bp else 0
        pct2 = (t_2 / t_bp * 100) if t_bp else 0
        self._txt(s, f"{pct0:.0f}% of best practices are NOT implemented. Only {pct2:.0f}% are fully implemented.",
                  0.8, 5.2, sz=16, bold=True, color=DB_ORANGE)

    # ====================== SLIDE 17: THANK YOU ======================

    def _s_thankyou(self, prs, host):
        s = prs.slides.add_slide(prs.slide_layouts[6]); self._bg(s)

        tx = s.shapes.add_textbox(Inches(2), Inches(1.5), Inches(9.5), Inches(1))
        p = tx.text_frame.paragraphs[0]
        p.text = "Thank You"; p.font.size = Pt(54); p.font.bold = True; p.font.color.rgb = DB_RED; p.alignment = PP_ALIGN.CENTER

        self._txt(s, f"Well-Architected Lakehouse Assessment \u2014 {host}", 2, 2.5, width=9.5, sz=14, color=DB_GRAY)

        self._txt(s, "Deliverables:", 0.8, 3.5, sz=14, bold=True)
        deliverables = [
            "This Executive Readout Deck (WAL_Assessment_Presentation.pptx)",
            "Detailed Assessment Report (WAL_Assessment_Readout.md)",
            "Scored Assessment Tool (WAL_Assessment_Scores.csv)",
            "Complete Audit Trail (WAL_Assessment_Audit_Report.md)",
        ]
        self._numbered(s, deliverables, 1.0, 3.9, width=10, sz=12)

        self._txt(s, "References:", 0.8, 5.5, sz=14, bold=True)
        refs = [
            "Databricks Well-Architected Framework (docs.databricks.com)",
            "Well-Architected Lakehouse Introduction (docs.databricks.com/lakehouse-architecture)",
            "Databricks Well-Architected Lakehouse Assessment Delivery Playbook (Internal)",
            "Databricks Well-Architected Lakehouse Assessment Tool (Internal)",
        ]
        self._bullets(s, refs, 1.0, 5.9, width=10, sz=10, color=DB_GRAY)
